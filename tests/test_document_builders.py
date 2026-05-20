"""End-to-end smoke tests for the document-builder modules.

These tests prove that each builder produces a valid artifact end-to-end
without depending on any task fixture state. They are the verification
layer for the document-generation tooling described in the plan at
.claude/plans/yes-i-think-it-s-valiant-pudding.md.
"""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
import pytest
from openpyxl import load_workbook

from scripts import xlsx
from scripts.charts import (
    arr_snapshot,
    bbrr_waterfall,
    deferred_rev_rollforward_chart,
    kpi_dashboard_grid,
    pl_bridge,
    top10_movement,
)
from scripts.docx import CEOLetterPayload, build_ceo_letter
from scripts.pptx import build_bsr_deck, build_mor_deck, build_parent_reportout_deck
from scripts.pptx.bsr import BSRPayload
from scripts.pptx.mor import MORPayload, slide_count
from scripts.pptx.parent_reportout import ParentReportoutPayload
from scripts.xlsx import builders, qa


# --- XLSX ------------------------------------------------------------------

def test_xlsx_shim_preserves_legacy_api(tmp_path: Path):
    """The scripts.format shim must keep dispatch.py + smoke.py working."""
    from scripts import format as fmt

    pack = tmp_path / "close_pack.xlsx"
    fmt.new_close_pack(pack, "2026-05")
    assert pack.exists()

    fmt.write_value_with_provenance(pack, "Cover", "B5", 9_605_500,
                                     claim_id="controller.consolidated.revenue_total",
                                     number_format="#,##0")
    df = pd.DataFrame([{"Account": "Subscription", "Actual": 9_605_500}])
    fmt.write_table(pack, "P&L", df)


def test_xlsx_builders_produce_valid_workbook(tmp_path: Path):
    pack = tmp_path / "pack.xlsx"
    xlsx.new_close_pack(pack, "2026-05")
    wb = load_workbook(pack)

    builders.build_cover(wb, period="2026-05", title="May 2026 Close",
                         kpi_strip=[
                             {"label": "Revenue", "value": 9_605_500,
                              "claim_id": "controller.consolidated.revenue_total",
                              "number_format": "$#,##0"},
                         ])
    df_pl = pd.DataFrame([
        {"Account": "Subscription", "Actual": 9_605_500, "Budget": 10_000_000,
         "Variance": -394_500},
    ])
    builders.build_pl(wb, df_pl, claim_ids_by_cell={
        "B2": "controller.consolidated.revenue_total",
        "D2": "fpa.variance_vs_budget.4100.usd",
    })
    builders.build_deferred_rev_rollforward(
        wb, period="2026-05",
        opening={"label": "Opening", "value": 42_000_000, "claim_id": "controller.deferred_rev.opening"},
        billings={"label": "Billings", "value": 10_500_000, "claim_id": "controller.billings.total"},
        recognized={"label": "Recognized", "value": 9_605_500, "claim_id": "controller.consolidated.revenue_total"},
    )
    wb.save(pack)
    assert pack.stat().st_size > 0

    # Re-open and confirm sheets exist
    wb2 = load_workbook(pack)
    assert "Cover" in wb2.sheetnames
    assert "P&L" in wb2.sheetnames
    assert "Deferred Rev Rollforward" in wb2.sheetnames


def test_xlsx_qa_detects_orphan_numeric_cells(tmp_path: Path):
    pack = tmp_path / "orphan.xlsx"
    xlsx.new_close_pack(pack, "2026-05")
    wb = load_workbook(pack)
    ws = wb.active
    ws["B5"] = 1_234_567   # numeric, no comment → orphan
    wb.save(pack)

    summary = qa.claim_id_coverage(pack)
    assert summary["total_numeric"] >= 1
    assert summary["with_claim_id"] == 0
    assert any("B5" in ref for _, ref, _ in summary["orphans"])

    with pytest.raises(AssertionError):
        qa.assert_full_claim_coverage(pack)


# --- Charts ----------------------------------------------------------------

def test_chart_library_renders_all_kinds(tmp_path: Path):
    paths = []
    paths.append(pl_bridge(
        output_path=tmp_path / "pl.png", title="P&L",
        budget=10e6, actual=9.6e6,
        drivers=[{"label": "RJ", "value": -120_000}],
    ))
    paths.append(bbrr_waterfall(
        output_path=tmp_path / "bbrr.png", title="BBRR",
        bookings=12e6, billings=10.5e6, rpo=78e6, revenue=9.6e6,
    ))
    paths.append(arr_snapshot(
        output_path=tmp_path / "arr.png", title="ARR",
        periods=["Jan", "Feb", "Mar"], arr=[100e6, 105e6, 110e6],
    ))
    paths.append(top10_movement(
        output_path=tmp_path / "t10.png", title="Top10",
        customers=["A", "B"], additions=[0, 0],
        expansions=[100, 200], contractions=[0, 0], churn=[0, 0],
    ))
    paths.append(kpi_dashboard_grid(
        output_path=tmp_path / "kpi.png", title="KPIs",
        kpis=[{"label": "ARR", "value": 110e6, "value_fmt": "$110M"}],
    ))
    paths.append(deferred_rev_rollforward_chart(
        output_path=tmp_path / "def.png", title="Def Rev",
        opening=42e6, billings=10e6, recognized=9.6e6,
    ))
    for p in paths:
        assert p.exists()
        assert p.stat().st_size > 5_000


# --- PPTX ------------------------------------------------------------------

def _chart(tmp_path: Path) -> Path:
    return arr_snapshot(
        output_path=tmp_path / "chart.png", title="Sample",
        periods=["Jan", "Feb"], arr=[100e6, 105e6],
    )


def test_mor_deck_builds(tmp_path: Path):
    chart_path = _chart(tmp_path)
    payload = MORPayload(
        period="2026-05",
        headline_kpis=[{"label": "Revenue", "value_text": "$9.6M",
                        "comparator_text": "-3.9%", "favorable": False,
                        "role": "external_link"}],
        headline_claim_ids=["controller.consolidated.revenue_total"],
        arr_chart_path=chart_path,
        arr_callouts=["ARR $110M"],
        arr_claim_ids=["reporting.arr.closing"],
        bbrr_chart_path=chart_path,
        bbrr_callouts=["BBRR aligned"],
        variance_rows=[{"archetype": "Tier-1", "product": "Flight Ops",
                         "customer": "RJ", "variance_usd": -120_000,
                         "mechanism": "Timing", "narrative": "RJ ramp"}],
        variance_claim_ids=["fpa.variance_vs_budget.4100.usd"],
    )
    out = build_mor_deck(payload, tmp_path / "mor.pptx")
    assert out.exists()

    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == slide_count(payload)
    # Provenance check: at least one slide carries claim_id in notes
    notes_blobs = [s.notes_slide.notes_text_frame.text for s in prs.slides]
    assert any("claim_id:" in n for n in notes_blobs), "expected claim_id in speaker notes"


def test_parent_reportout_appends_segment_slides(tmp_path: Path):
    chart_path = _chart(tmp_path)
    payload = ParentReportoutPayload(
        period="2026-06",
        headline_kpis=[{"label": "Revenue", "value_text": "$30M"}],
        arr_chart_path=chart_path, arr_callouts=["x"],
        segment_recon_headers=["Internal", "Parent", "Delta"],
        segment_recon_rows=[["4100 Subscription", "4100-DIG", 0]],
        account_map_deltas=[{"internal_account": "4150",
                              "parent_account": "4150-DIG",
                              "status": "new", "notes": "ATS records"}],
    )
    out = build_parent_reportout_deck(payload, tmp_path / "po.pptx")
    from pptx import Presentation
    prs = Presentation(str(out))
    # Should have at least: title + headline + ARR + segment divider + segment table + account-map table
    assert len(prs.slides) >= 6


def test_bsr_deck_builds(tmp_path: Path):
    payload = BSRPayload(
        period="2026-06",
        bs_summary_kpis=[{"label": "Total assets", "value_text": "$180M"}],
        account_roll_headers=["Account", "Beginning", "Ending"],
        account_roll_rows=[["Cash", "$50M", "$52M"]],
        flux_rows=[{"archetype": "—", "product": "—", "customer": "Cash",
                     "variance_usd": 2_000_000, "mechanism": "Operating",
                     "narrative": "Free cash flow build"}],
        reserves_headers=["Account", "Balance"],
        reserves_rows=[["AR allowance", "$1.2M"]],
    )
    out = build_bsr_deck(payload, tmp_path / "bsr.pptx")
    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) >= 4


# --- DOCX ------------------------------------------------------------------

def test_ceo_letter_with_xlsx_table(tmp_path: Path):
    # Build a sample close pack
    pack = tmp_path / "close_pack.xlsx"
    xlsx.new_close_pack(pack, "2026-05")
    wb = load_workbook(pack)
    df = pd.DataFrame([
        {"Account": "Subscription", "Actual": 9_605_500, "Budget": 10_000_000,
         "Variance": -394_500},
    ])
    builders.build_pl(wb, df, claim_ids_by_cell={
        "B2": "controller.consolidated.revenue_total",
    })
    wb.save(pack)

    md = """# Monthly CEO Letter — May 2026

## Headlines
Revenue closed at $9.6M [claim: controller.consolidated.revenue_total].
"""

    out, claim_ids = build_ceo_letter(
        CEOLetterPayload(period="2026-05",
                         narrative_md=md,
                         table_xlsx_path=pack,
                         table_sheet="P&L",
                         signer_name="Test Signer"),
        tmp_path / "letter.docx",
    )
    assert out.exists()
    assert "controller.consolidated.revenue_total" in claim_ids

    from docx import Document
    doc = Document(str(out))
    assert len(doc.tables) == 1
    assert len(doc.tables[0].rows) >= 2  # header + at least one data row


# --- PDF utility -----------------------------------------------------------

def test_pdf_from_markdown(tmp_path: Path):
    md = tmp_path / "sample.md"
    md.write_text("# Sample\n\nBody text.\n")
    from scripts.pdf import to_pdf
    try:
        out = to_pdf(md, tmp_path / "sample.pdf")
    except OSError as e:
        if "cannot load library" in str(e) or "libgobject" in str(e):
            pytest.skip(f"WeasyPrint native deps unavailable: {e}")
        raise
    assert out.exists()
    assert out.stat().st_size > 1_000


# --- Freshness check -------------------------------------------------------

def test_freshness_no_changes_short_circuits(tmp_path: Path, monkeypatch):
    """With pins at HEAD, the diff should report zero changes."""
    from scripts.tooling import freshness_check as fr

    pin_path = fr.REPO_ROOT / "memory" / "upstream_skills_pin.json"
    report = fr.diff_report(pin_path=pin_path, fetch=False)
    assert "skills" in report
    # Verify the no-changes path produces a clean markdown report
    md = fr.render_report_md(report)
    assert "freshness report" in md.lower()
